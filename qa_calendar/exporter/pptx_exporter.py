"""
QA 캘린더 PPTX 내보내기
계획안 PPTX 템플릿을 그대로 복제하여 내용을 채움
"""
import copy
import re
from pathlib import Path
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor

# 계획안 템플릿 경로
TEMPLATE_PATH = Path(__file__).parent.parent.parent / 'templates' / 'calendar_template.pptx'

MONTH_KEYS = ['01', '02', '03', '04', '05', '06',
              '07', '08', '09', '10', '11', '12']


def export_to_pptx(output_path: str, calendar_data: dict, issues_data: dict):
    """
    캘린더 데이터를 PPTX로 내보내기.
    템플릿 PPTX 기반으로 편집 가능 셀 내용 교체.

    calendar_data: { slide1: { rows: [[{text, rowspan, colspan, editable}, ...], ...] } }
    issues_data:   { slide2: { rows: [...] } }
    """
    if not TEMPLATE_PATH.exists():
        raise FileNotFoundError(f"템플릿 없음: {TEMPLATE_PATH}")

    prs = Presentation(str(TEMPLATE_PATH))

    slide1_rows = calendar_data.get('slide1', {}).get('rows', [])
    slide2_rows = issues_data.get('slide2', {}).get('rows', [])

    if slide1_rows:
        _fill_table(prs.slides[0], slide1_rows)
    if slide2_rows:
        _fill_table(prs.slides[1], slide2_rows)

    prs.save(output_path)


def _fill_table(slide, rows_data: list):
    """슬라이드에서 테이블을 찾아 편집 가능 셀 내용 교체"""
    for shape in slide.shapes:
        if not hasattr(shape, 'table'):
            continue
        tbl = shape.table
        tbl_rows = list(tbl.rows)

        # 물리적 셀 위치 계산 (병합 고려)
        # pptx는 병합된 셀도 실제 _tc 요소를 가지고 있으나
        # cells[ci]는 실제 물리 셀에 접근
        for ri, row_data in enumerate(rows_data):
            if ri >= len(tbl_rows):
                break

            phys_cells = list(tbl_rows[ri].cells)
            ci = 0  # 물리 셀 인덱스

            for cell_data in row_data:
                if ci >= len(phys_cells):
                    break

                colspan = cell_data.get('colspan', 1)
                text = cell_data.get('text', '')
                editable = cell_data.get('editable', False)

                if editable and text is not None:
                    _write_cell(phys_cells[ci], str(text))

                ci += max(colspan, 1)

        break  # 첫 번째 테이블만 처리


def _write_cell(cell, text: str):
    """셀에 텍스트 쓰기 (기존 단락 서식 유지)"""
    tf = cell.text_frame
    if not tf:
        return

    # 기존 서식 참조 (첫 단락 첫 런)
    ref_size = None
    ref_bold = None
    ref_name = None
    ref_color = None

    for para in tf.paragraphs:
        for run in para.runs:
            if run.font.size:
                ref_size = run.font.size
            if run.font.bold is not None:
                ref_bold = run.font.bold
            if run.font.name:
                ref_name = run.font.name
            try:
                if run.font.color.type is not None:
                    ref_color = run.font.color.rgb
            except Exception:
                pass
            if ref_size:
                break
        if ref_size:
            break

    # 기존 단락 내용 클리어 (첫 단락 유지, 나머지 삭제)
    # lxml로 직접 조작
    from lxml import etree
    txBody = tf._txBody

    # 기존 <a:p> 요소들 수집
    nsmap = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    paras = txBody.findall(f'{{{nsmap}}}p')

    # 첫 단락 서식 보존용
    first_para = paras[0] if paras else None

    # 모든 단락 제거
    for p in paras:
        txBody.remove(p)

    # 새 내용 작성
    lines = text.split('\n') if text else ['']

    for i, line in enumerate(lines):
        # 새 단락 생성
        p_el = etree.SubElement(txBody, f'{{{nsmap}}}p')

        # 첫 단락의 pPr(단락 속성) 복사
        if first_para is not None:
            pPr = first_para.find(f'{{{nsmap}}}pPr')
            if pPr is not None:
                p_el.insert(0, copy.deepcopy(pPr))

        # 런 생성
        r_el = etree.SubElement(p_el, f'{{{nsmap}}}r')

        # 런 속성
        rPr = etree.SubElement(r_el, f'{{{nsmap}}}rPr', attrib={'lang': 'ko-KR', 'altLang': 'en-US'})

        if ref_size:
            rPr.set('sz', str(int(ref_size.pt * 100)))
        else:
            rPr.set('sz', '900')  # 기본 9pt

        if ref_name:
            latin = etree.SubElement(rPr, f'{{{nsmap}}}latin', attrib={'typeface': ref_name})

        if ref_color:
            solidFill = etree.SubElement(rPr, f'{{{nsmap}}}solidFill')
            srgbClr = etree.SubElement(solidFill, f'{{{nsmap}}}srgbClr',
                                        attrib={'val': str(ref_color)})

        # 텍스트 요소
        t_el = etree.SubElement(r_el, f'{{{nsmap}}}t')
        t_el.text = line

    # txBody 끝에 <a:lstStyle/> 없으면 단락들 뒤에 위치하도록 정렬
    # (pptx 규격상 lstStyle은 p 앞에 있어야 함 - 이미 있을 것)
