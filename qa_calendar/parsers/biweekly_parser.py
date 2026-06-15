"""
바이위클리 PPTX 파서

슬라이드 구조:
- Ⅰ. 기획점검 및 모니터링 슬라이드들
- Ⅱ. 상품QA 진행 / 이슈 슬라이드들
  - 사전QA 진행 이슈상품 (방송상품)
  - 현장QA 이슈상품
  - 암행점검 부적합 상품
"""
import re
from pptx import Presentation
from pptx.util import Pt
from dataclasses import dataclass
from typing import Optional


@dataclass
class QAIssue:
    category: str       # '방송상품' | '현장QA' | '암행'
    product_name: str
    md_name: str
    issue: str
    measure: str
    result: str         # 개선 | 보류 | 부적합 (현장)


@dataclass
class BiweeklyData:
    date: str           # "2026. 3. 25."
    month: str          # '03'
    qa_issues: list     # List[QAIssue]
    inspection_titles: list  # 기획점검/모니터링 항목 제목 리스트


def parse_biweekly(filepath: str) -> BiweeklyData:
    """바이위클리 PPTX 파싱"""
    prs = Presentation(filepath)

    date_str = ''
    month = ''
    qa_issues = []
    inspection_titles = []

    for slide in prs.slides:
        slide_text = _get_slide_text(slide)
        slide_title = _get_slide_title(slide_text)

        # 날짜 추출 (첫 슬라이드)
        if not date_str:
            m = re.search(r'2026\.\s*(\d+)\.\s*\d+\.', slide_text)
            if m:
                date_str = m.group(0).strip()
                month = m.group(1).zfill(2)

        # Ⅰ. 기획점검/모니터링 슬라이드
        if 'Ⅰ' in slide_title or '기획점검' in slide_title:
            titles = _extract_inspection_titles(slide)
            inspection_titles.extend(titles)

        # Ⅱ. 상품QA 진행이슈 슬라이드 (제목 또는 본문에 관련 키워드 포함)
        if ('Ⅱ' in slide_text or 'QA' in slide_title or '이슈' in slide_title
                or '암행점검' in slide_text or '사전QA' in slide_text
                or '현장QA' in slide_text or '현장 QA' in slide_text):
            issues = _extract_qa_issues(slide, slide_text)
            qa_issues.extend(issues)

    return BiweeklyData(
        date=date_str,
        month=month,
        qa_issues=qa_issues,
        inspection_titles=inspection_titles,
    )


def _get_slide_text(slide) -> str:
    """슬라이드 전체 텍스트 추출"""
    texts = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            texts.append(shape.text_frame.text)
    return '\n'.join(texts)


def _get_slide_title(slide_text: str) -> str:
    """슬라이드 제목 줄 추출"""
    for line in slide_text.split('\n'):
        line = line.strip()
        if line and len(line) > 3:
            return line
    return ''


def _extract_inspection_titles(slide) -> list:
    """기획점검/모니터링 슬라이드에서 항목 제목 추출"""
    titles = []
    for shape in slide.shapes:
        if not hasattr(shape, 'table'):
            continue
        tbl = shape.table
        rows_list = list(tbl.rows)
        if not rows_list:
            continue
        # 헤더 확인
        header = [c.text_frame.text.strip() for c in rows_list[0].cells]
        if 'Title' in header or '결과' in header:
            title_col = 0
            for i, h in enumerate(header):
                if h == 'Title' or i == 0:
                    title_col = i
                    break
            for row in rows_list[1:]:
                cells = [c.text_frame.text.strip() for c in row.cells]
                if cells and cells[title_col]:
                    title = cells[title_col].replace('\n', ' ').strip()
                    if title and title not in titles:
                        titles.append(title)
    return titles


def _extract_qa_issues(slide, slide_text: str) -> list:
    """상품QA 이슈 슬라이드에서 이슈 추출"""
    issues = []

    # 카테고리 파악
    category = _detect_category(slide_text)
    if not category:
        return issues

    for shape in slide.shapes:
        if not hasattr(shape, 'table'):
            continue
        tbl = shape.table
        qa_rows = list(tbl.rows)
        if not qa_rows:
            continue

        headers = [c.text_frame.text.strip() for c in qa_rows[0].cells]
        if '상품명' not in headers[0] and '상품명 (MD)' not in headers[0]:
            continue

        for row in qa_rows[1:]:
            cells = [c.text_frame.text.strip() for c in row.cells]
            if not cells or not cells[0]:
                continue

            product_raw = cells[0]
            issue_raw = cells[1] if len(cells) > 1 else ''

            product_name, md_name, sourcing = _parse_product_cell(product_raw)
            issue_text, measure_text, result = _parse_issue_cell(issue_raw, category)

            if product_name:
                issues.append(QAIssue(
                    category=category,
                    product_name=product_name,
                    md_name=md_name,
                    issue=issue_text,
                    measure=measure_text,
                    result=result,
                ))

    return issues


def _detect_category(slide_text: str) -> Optional[str]:
    """슬라이드 텍스트로 QA 이슈 카테고리 감지"""
    if '사전QA' in slide_text or '사전 QA' in slide_text:
        return '방송상품'
    if '현장QA' in slide_text or '현장 QA' in slide_text:
        return '현장QA'
    if '암행점검' in slide_text or '암행 점검' in slide_text:
        return '암행'
    return None


def _parse_product_cell(text: str) -> tuple:
    """상품명(MD) 셀 파싱 → (상품명, MD명, 소싱매체)"""
    sourcing = ''
    if '[TV]' in text:
        sourcing = 'TV'
    elif '[모바일]' in text:
        sourcing = '모바일'

    # 소싱 태그 제거
    text = re.sub(r'\[(TV|모바일|라이브)\]\s*', '', text).strip()

    # MD명 추출 (괄호 안 "xxx M")
    md_match = re.search(r'\(([^)]+M(?:/[^)]+M)?)\)', text)
    md_name = md_match.group(1) if md_match else ''

    # MD 괄호 제거 후 상품명
    product = re.sub(r'\([^)]+M(?:/[^)]+M)?\)', '', text).strip()
    product = product.replace('\n', ' ').strip()

    return product, md_name, sourcing


def _parse_issue_cell(text: str, category: str) -> tuple:
    """이슈사항 셀 파싱 → (이슈내용, 조치사항, 결과)"""
    issue = ''
    measure = ''
    result = ''

    if category == '현장QA':
        # [상세내용] 섹션
        issue_match = re.search(r'\[상세내용\](.*?)(?:\[조치사항\]|$)', text, re.DOTALL)
        measure_match = re.search(r'\[조치사항\](.*?)$', text, re.DOTALL)
        result_match = re.search(r'(부적합|보류|적합)', text)

        issue = issue_match.group(1).strip() if issue_match else text.strip()
        measure = measure_match.group(1).strip() if measure_match else ''
        result = result_match.group(1) if result_match else ''

    elif category in ('방송상품', '암행'):
        # [이슈사항] or [부적합 사항] 섹션
        issue_match = re.search(r'\[(?:이슈사항|부적합 사항)\](.*?)(?:\[조치사항\]|$)', text, re.DOTALL)
        measure_match = re.search(r'\[조치사항\](.*?)$', text, re.DOTALL)
        result_match = re.search(r'상품QA\s*[–-]\s*(개선|보류|부적합)', text)

        issue = issue_match.group(1).strip() if issue_match else text.strip()
        measure = measure_match.group(1).strip() if measure_match else ''
        result = result_match.group(1) if result_match else ''

    # 줄바꿈 정리
    issue = re.sub(r'\n+', '\n', issue).strip()[:300]
    measure = re.sub(r'\n+', ' ', measure).strip()[:150]

    return issue, measure, result
